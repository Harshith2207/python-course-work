from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

prs = Presentation()

def add_story_slide(title, story_lines, code_snippet, image_path=None):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    
    # Add optional image
    if image_path:
        try:
            slide.shapes.add_picture(image_path, Inches(0.3), Inches(1.5), height=Inches(3))
        except FileNotFoundError:
            pass  # Skip if no image found

    # Add story text
    text_box = slide.shapes.add_textbox(Inches(3.8), Inches(1.2), Inches(5), Inches(3))
    tf = text_box.text_frame
    tf.word_wrap = True
    for line in story_lines:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.font.name = "Segoe UI"
    
    # Add code box
    if code_snippet:
        code_box = slide.shapes.add_textbox(Inches(3.8), Inches(3.5), Inches(5), Inches(2))
        tf_code = code_box.text_frame
        run = tf_code.paragraphs[0].add_run()
        run.text = code_snippet
        run.font.name = 'Consolas'
        run.font.size = Pt(14)
        run.font.color.rgb = RGBColor(255, 255, 255)
        code_box.fill.solid()
        code_box.fill.fore_color.rgb = RGBColor(30, 30, 30)

# Slide 1 – Title
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Lily & Py the Robot Learn Python Functions"
slide.placeholders[1].text = "A Kid-Friendly Story about User-Defined Functions"

# Slide 2 – What is a Function
add_story_slide(
    "What is a Function?",
    ["Lily has a magic robot, Py.", 
     "Py can do anything, but first Lily must teach it a trick.",
     "In Python, we teach tricks using 'functions'."],
    """def say_hello():
    print("Hello!")

say_hello()  # Py says Hello!""",
    "robot.png"
)

# Slide 3 – With & Without Parameters
add_story_slide(
    "With & Without Parameters",
    ["Sometimes Py needs information to work (parameters).",
     "Sometimes he just acts with no input."],
    """def greet(name):
    print("Hi", name)

def wave():
    print("*waves*")"""
)

# Slide 4 – With & Without Return Values
add_story_slide(
    "With & Without Return Values",
    ["Sometimes Py gives back an answer (return).",
     "Sometimes he just does something without replying."],
    """def add(a, b):
    return a + b

def sing():
    print("La la la!")"""
)

# Slide 5 – Default Parameters & Summary
add_story_slide(
    "Default Parameters & Summary",
    ["Py has default tricks if you don't tell him exactly what to do.",
     "Functions make coding easy, reusable, and fun!"],
    """def greet(name="Friend"):
    print("Hello", name)"""
)

prs.save("Lily_and_Py_Functions.pptx")
print("✅ 5-slide kids-style PPT created: Lily_and_Py_Functions.pptx")
